from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import os

def generate_presentation(vacancies, background_path, output_path):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for v in vacancies:
        slide = prs.slides.add_slide(blank_slide_layout)

        # Установка фона слайда
        if os.path.exists(background_path):
            slide.shapes.add_picture(background_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(5)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_bottom = Inches(0.1)

        def add_paragraph(text, bold=False, size=Pt(28), color=(0, 0, 0), link=None):
            p = tf.add_paragraph()
            p.text = text
            p.font.bold = bold
            p.font.size = size
            p.font.color.rgb = RGBColor(*color)
            p.alignment = PP_ALIGN.LEFT
            if link:
                run = p.runs[0]
                run.hyperlink.address = link

        tf.clear()
        add_paragraph(f"📌 {v['title']}", bold=True, size=Pt(36))
        add_paragraph(f"💰 {v['salary']}")
        add_paragraph(f"📍 {v['location']}")
        add_paragraph(f"🏢 {v['employer']}", link=v.get('employer_link'))
        add_paragraph(f"⚙️ {v['conditions']}")
        add_paragraph(f"📞 {v['contact']}", link=v.get('contact'))
        add_paragraph(f"🌐 Источник", link=v.get('source'))

    prs.save(output_path)
